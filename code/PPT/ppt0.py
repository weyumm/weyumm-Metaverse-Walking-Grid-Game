from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.util import Pt

# 创建PPT对象
prs = Presentation()

# 添加背景图片的方法
def set_slide_background(slide, image_path):
    left = top = 0
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    background = slide.shapes.add_picture(image_path, left, top, slide_width, slide_height)
    slide.shapes._spTree.remove(background._element)
    slide.shapes._spTree.insert(2, background._element)

# 添加白色文本框
def add_white_textbox(slide, title_text, content_text):
    title = slide.shapes.title
    title.text = title_text
    for paragraph in title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(255, 255, 255)

    textbox = slide.shapes.placeholders[1]
    textbox.text = content_text
    for paragraph in textbox.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(255, 255, 255)
            
# 设置字体稍小
def add_white_textbox_with_smaller_font(slide, title_text, content_text):
    title = slide.shapes.title
    title.text = title_text
    for paragraph in title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(28)  # 标题字体略小
            run.font.color.rgb = RGBColor(255, 255, 255)

    textbox = slide.shapes.placeholders[1]
    textbox.text = content_text
    for paragraph in textbox.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(18)  # 内容字体小一点
            run.font.color.rgb = RGBColor(255, 255, 255)

# 封面
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide, 'PPT背景.jpg')
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "元宇宙视角下开放世界探索游戏的开发前景——以米哈游的游戏产品为例"
subtitle.text = "基于米哈游的《崩坏3》《原神》《崩坏：星穹铁道》《绝区零》"
subtitle.text = "张恒祯\n同济大学汽车院"

# 幻灯片1: 什么是元宇宙
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide, 'PPT背景.jpg')
add_white_textbox(slide, "什么是元宇宙？", 
                  "元宇宙是一个虚拟现实、增强现实和区块链融合的数字生态系统，它通过高度沉浸的虚拟世界将玩家带入全新体验。"
                  "\n\n特点：\n- 高自由度\n- 虚拟经济系统\n- 社交和跨平台沉浸体验")

# 幻灯片2: 崩坏3 分析
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide, 'PPT背景.jpg')
add_white_textbox(slide, "游戏概述：崩坏3", 
                  "《崩坏3》作为米哈游开创性作品，以多维度的虚拟世界和叙事为特色。"
                  "\n\n元宇宙特点：\n- 虚拟世界与现实联动\n- 高互动性与沉浸式战斗体验\n"
                  "- 角色与世界的深度联系，通过“作战小队”探索未知领域。")

# 幻灯片4: 元宇宙视角下的崩坏3体验
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide, 'PPT背景.jpg')

add_white_textbox_with_smaller_font(slide, "元宇宙视角下的《崩坏3》体验", 
    "《崩坏3》为玩家带来了极具沉浸感的虚拟战斗体验，结合了多维度叙事和多样化的战斗机制。"
    "\n\n体验特点：\n- 玩家能够通过多种角色进行互动，沉浸在剧情发展的虚拟世界中。"
    "\n- 游戏中的虚拟小队战斗系统让玩家感受到真实的团队作战体验。"
    "\n- 多样化的任务和关卡设计让游戏充满变化与挑战，进一步增强了元宇宙的互动性。"
    "\n\n进步意义：\n- 《崩坏3》通过引入丰富的叙事元素和角色成长体系，证明了元宇宙中故事与世界构建的深度结合。"
    "\n- 游戏的多人在线战斗模式展现了虚拟社交与战斗的深层互动，推动了元宇宙中的虚拟合作体验。"
    "\n- 《崩坏3》在全球范围内的成功进一步推动了虚拟经济和虚拟角色的成长，标志着元宇宙的早期探索和发展。")


# 幻灯片3: 原神 分析
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide, 'PPT背景.jpg')
add_white_textbox(slide, "游戏分析：原神", 
                  "《原神》是开放世界的代表作，其极高自由度和丰富的探索元素被认为是元宇宙游戏的重要前驱之一。"
                  "\n\n元宇宙特点：\n- 无缝开放世界探索，玩家可以随意探索解谜\n- 虚拟经济系统与角色交易\n"
                  "- 多维度社交互动，虚拟世界和现实交融。")

# 幻灯片3: 元宇宙视角下的原神体验
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide, 'PPT背景.jpg')

add_white_textbox_with_smaller_font(slide, "元宇宙视角下的《原神》体验", 
    "《原神》作为米哈游的开放世界游戏，为玩家提供了广阔的虚拟探索空间，深度结合了元宇宙的核心理念。"
    "\n\n体验特点：\n- 开放的地图探索与自由的解谜机制让玩家感受到无缝衔接的虚拟世界。"
    "\n- 在多人协作模式中，玩家可以与全球其他用户共同探索，形成真正的虚拟社交互动。"
    "\n- 角色养成、元素反应与战斗机制为元宇宙中的沉浸式战斗提供了创新体验。"
    "\n\n进步意义：\n- 《原神》证明了虚拟与现实的边界可以通过技术模糊，提供无缝虚拟体验。"
    "\n- 跨平台支持（PC、移动设备、主机）打破了设备壁垒，推动了元宇宙游戏的普及。"
    "\n- 游戏内的虚拟经济系统和不断扩展的内容更新证明了可持续发展的游戏生态。")


# 幻灯片4: 崩坏：星穹铁道 分析
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide, 'PPT背景.jpg')
add_white_textbox(slide, "游戏分析：崩坏：星穹铁道", 
                  "《崩坏：星穹铁道》融合了战斗与探索，玩家可以穿越不同的星球和次元，带来多维度体验。"
                  "\n\n元宇宙特点：\n- 多维世界探索，跨越星球和次元，符合作为元宇宙游戏的核心理念。\n"
                  "- 虚拟角色和虚拟物品交易系统，与虚拟现实深度整合。")

# 幻灯片5: 元宇宙视角下的《崩坏：星穹铁道》体验
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide, 'PPT背景.jpg')

add_white_textbox_with_smaller_font(slide, "元宇宙视角下的《崩坏：星穹铁道》体验", 
    "《崩坏：星穹铁道》带来了跨越星球和维度的多维探索体验，"
    "玩家通过乘坐星际列车穿梭于不同星球之间，展开冒险与战斗。"
    "\n\n体验特点：\n- 多维度的星际探索带来了广阔的元宇宙体验，"
    "每个星球都有独特的文化、生态与挑战。"
    "\n- 游戏引入了策略性回合制战斗，结合角色养成，使玩家更沉浸在虚拟世界的复杂互动中。"
    "\n- 虚拟团队合作与任务系统，体现了元宇宙中的社会互动和跨星球协作。"
    "\n\n进步意义：\n- 《崩坏：星穹铁道》展示了元宇宙游戏的可能性，即通过跨维度的虚拟世界构建，"
    "增强了虚拟空间的广度与深度。"
    "\n- 游戏利用虚拟交通工具（星际列车）实现了元宇宙中的无缝移动，提升了玩家在虚拟世界中的沉浸感。"
    "\n- 《崩坏：星穹铁道》通过策略性战斗和开放式探索，扩展了玩家在元宇宙中的互动方式。")


# 幻灯片5: 绝区零 分析
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide, 'PPT背景.jpg')
add_white_textbox(slide, "游戏分析：绝区零", 
                  "《绝区零》带来了全新的战斗体验和末世氛围，结合AR技术，为玩家提供了紧凑的游戏节奏。"
                  "\n\n元宇宙特点：\n- 高度沉浸的战斗系统和虚拟世界交互\n- 增强现实元素与虚拟世界融合。")
# 幻灯片6: 元宇宙视角下的《绝区零》体验
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide, 'PPT背景.jpg')

add_white_textbox_with_smaller_font(slide, "元宇宙视角下的《绝区零》体验", 
    "《绝区零》带来了紧张的战斗节奏和末世氛围，结合了高度自由的探索与战斗机制。"
    "\n\n体验特点：\n- 游戏中的虚拟现实与增强现实元素，模糊了虚拟与现实的边界，"
    "为玩家提供了更加沉浸的元宇宙体验。"
    "\n- 战斗系统节奏紧凑，强调虚拟与现实的融合，玩家可以通过即时反应完成复杂的战斗操作。"
    "\n- 通过虚拟任务和虚拟环境交互，体现了元宇宙中的开放式冒险体验。"
    "\n\n进步意义：\n- 《绝区零》通过虚拟与现实的深度融合，探索了虚拟世界与现实世界之间的互动边界，"
    "为元宇宙的虚拟交互提供了新的范例。"
    "\n- 游戏引入了复杂的战斗机制和末世背景，增强了虚拟世界的情感张力和玩家的代入感。"
    "\n- 《绝区零》展示了元宇宙中沉浸式虚拟现实技术与游戏设计的高度融合，为未来元宇宙游戏奠定了基础。")


# 幻灯片6: 元宇宙游戏中的技术支持
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide, 'PPT背景.jpg')
add_white_textbox(slide, "元宇宙游戏中的技术支持", 
                  "1. 虚拟现实技术：提供沉浸式的体验，带领玩家进入虚拟世界。\n"
                  "2. 区块链技术：确保虚拟资产的独立性和安全性，通过NFT支持虚拟物品的交易。\n"
                  "3. 云计算与5G：大规模多人在线互动，支持实时数据同步和动态内容生成。")

# 幻灯片7: 生成式人工智能对元宇宙游戏开发的帮助
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide, 'PPT背景.jpg')
add_white_textbox(slide, "生成式人工智能对元宇宙游戏开发的帮助", 
                  "生成式AI通过自动生成角色、场景、任务等，大大减少了开发工作负担。\n"
                  "走格子游戏Demo展示：\n- AI生成虚拟场景与角色，展示生成式AI在游戏设计中的应用。\n"
                  "- 在元宇宙游戏中，AI不仅可以生成内容，还可以个性化游戏体验，带来无限可能。")

# 幻灯片8: 结论
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide, 'PPT背景.jpg')
add_white_textbox(slide, "结论", 
                  "米哈游的几款游戏充分展示了元宇宙游戏产品的潜力，"
                  "它们结合了虚拟现实技术和增强现实技术，创造了高沉浸的虚拟世界。\n"
                  "生成式人工智能为元宇宙游戏开发带来了巨大的可能性，未来将成为重要推动力量。")


# 保存PPT文件
prs.save('元宇宙游戏产品分析.pptx')
